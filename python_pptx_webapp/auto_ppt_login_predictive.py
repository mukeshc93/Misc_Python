from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR,PP_ALIGN
from pptx.chart import *
from pptx.enum.chart import XL_CHART_TYPE,XL_LEGEND_POSITION,XL_TICK_MARK,XL_TICK_LABEL_POSITION
from pptx.chart.data import ChartData
import pandas as pd
import numpy as np
from os import listdir
from os.path import isfile, join
import sys
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from StringIO import StringIO
import calendar

from flask import Flask, render_template_string, request,abort , redirect , Response ,url_for
from wtforms import Form, SelectMultipleField,SelectField
from flask import send_file,send_from_directory
from flask_login import LoginManager , login_required , UserMixin , login_user
import flask_login



app = Flask(__name__)

##Keep this file as the customer map file. Read it in the UI

usernames=pd.read_csv("/home/ubuntu/sdmauto/Names/Names.csv",error_bad_lines=False,encoding='utf-8')

##Accumulating Reporttype and Date period
mypath = "/home/ubuntu/sdmauto/reportdata"
onlyfiles = pd.DataFrame([f for f in listdir(mypath) if isfile(join(mypath, f))],columns=['Name'])
onlyfiles['Reporttype'] = onlyfiles.Name.str.split('.').str.get(0)


class partnerform(Form):
    #For multiple fields use SelectMultipleFields
    NAME = SelectField(u'Customer Name', choices=[(u, u) for u in sorted(usernames["COMPANY_NAME"].dropna().drop_duplicates().sort_values(),key=unicode.lower)])
    Reporttype = SelectMultipleField(u'Report Type ', choices=[(u, u) for u in sorted(onlyfiles["Reporttype"].dropna().drop_duplicates().sort_values())])

template_form = """
{% block content %}
<h1>Select Customer & Report details</h1>
<form method="POST" action="/sdm/report">
    <h2><div>{{ form.NAME.label }}<Label></div></h2>
    <div>{{ form.NAME(size=10) }}</div>
    <h2><div>{{ form.Reporttype.label }}<Label></div></h2>
    <div>{{ form.Reporttype(multiple=True,size=10) }}</div>
    <button type="submit" class="btn">Submit</button>    
</form>
{% endblock %}
"""

completed_template = """
{% block body %}
<h1>Select Customer & Report details</h1>
{% endblock %}
"""

app.config['SECRET_KEY'] = ''
login_manager = LoginManager()
#login_manager.login_view = "login"
login_manager.init_app(app)

users = {'admin@sdm': {'password': ''}}

class User(UserMixin):
    pass

@login_manager.user_loader
def user_loader(email):
    if email not in users:
        return

    user = User()
    user.id = email
    return user


@login_manager.request_loader
def request_loader(request):
    email = request.form.get('email')
    if email not in users:
        return

    user = User()
    user.id = email

    # DO NOT ever store passwords in plaintext and always compare password
    # hashes using constant-time comparison!
    user.is_authenticated = request.form['password'] == users[email]['password']

    return user



@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'GET':
        return '''
               <form action='login' method='POST'>
                <input type='text' name='email' id='email' placeholder='email'/>
                <input type='password' name='password' id='password' placeholder='password'/>
                <input type='submit' name='submit' value='login'/>
               </form>
               '''

    email = request.form['email']
    if request.form['password'] == users[email]['password']:
        user = User()
        user.id = email
        login_user(user)
        return redirect(url_for('home'))

    return 'Bad login'

@app.route('/logout')
def logout():
    flask_login.logout_user()
    return 'Logged out'

@login_manager.unauthorized_handler
def unauthorized_handler():
    return 'Unauthorized'

@app.route('/home')
@login_required
def home():
    form = partnerform(request.form)
    return render_template_string(template_form,form = form)

@app.route('/report',methods=['POST'])
def index():
    #print >>sys.stderr, 'Apache Nonsense!'
    form = partnerform(request.form)
    print "POST request and form is valid"
    name =  form.NAME.data
    print "Partner names %s" % name
    print >>sys.stderr, 'Partner names!'
    ##FIRST SLIDE
    prs = Presentation("/home/ubuntu/data/ppt/ppt.pptx")
    slide = prs.slides[0]
    shapes = slide.shapes    
    shapes[0].text = name.encode('utf8') + ' Monthly Review'
    shapes[1].text = 'Mukesh Choudhary \n' + calendar.month_name[pd.datetime.now().month] + ' ' + str(pd.datetime.now().year)

    empty_slide = prs.slide_masters[1].slide_layouts[5]
    slide = prs.slides.add_slide(empty_slide)
    shapes = slide.shapes

    ##Grey text box
    left = Inches(0.29) # 0.93" centers this overall set of shapes
    top = Inches(0.58)
    width = Inches(4.55)
    height = Inches(6.92)
    # s = shapes3.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s = shapes.add_textbox(left, top, width, height)
    fill=s.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xEF, 0xF0, 0xF1)

    #Orange textbox with title
    left = Inches(0) # 0.93" centers this overall set of shapes
    top = Inches(2.82)
    width = Inches(10)
    height = Inches(1.82)
    # s = shapes3.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s = shapes.add_textbox(left, top, width, height)
    fill=s.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xC0, 0x00)
    # s.text ='Step 1'
    tframe = s.text_frame
    # tframe.clear()
    tframe.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tframe.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.level=1
    run = p.add_run()
    run.text = 'EXECUTIVE SUMMARY'
    font = run.font
    font.name = 'Arial'
    font.size = Pt(38)
    font.bold = True 
    font.color.rgb = RGBColor(0x00, 0x00, 0x00) 

    #Footer for customer report
    left = Inches(6.63)
    top = Inches(6.84)
    width = Inches(3.26)
    height = Inches(0.4)
    s = shapes.add_textbox(left, top, width, height)
    tframe = s.text_frame
    tframe.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tframe.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = 'Customer Report'
    font = run.font
    font.name = 'Arial'
    font.size = Pt(8.2)
    font.color.rgb = RGBColor(127, 127, 127)
      
    for i in form.Reporttype.data:
        if "DAMAGE" in str(i):
		##FIRST CHART
		path = "/home/ubuntu/sdmauto/reportdata/" + str(i) +".csv"
		#print path
		df=pd.read_csv(path)
		#name='Consum'
		df = df.loc[df['COMPANYNAME'] == name.encode('utf8')][-12:].reset_index(drop=True)
		#df = df.groupby(['NAME','DATE'], as_index=False).sum()
		#print df.head(100)
		empty_slide = prs.slide_masters[1].slide_layouts[5]
		slide = prs.slides.add_slide(empty_slide)
		shapes = slide.shapes
		left = Inches(0) # 0.93" centers this overall set of shapes
		top = Inches(0.8)
		width = Inches(10)
		height = Inches(0.4)
		# s = shapes3.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
		s = shapes.add_textbox(left, top, width, height)
		fill=s.fill
		fill.solid()
		fill.fore_color.rgb = RGBColor(0xFF, 0xC0, 0x00)
		# s.text ='Step 1'
		tframe = s.text_frame
		# tframe.clear()
		p = tframe.paragraphs[0]
		p.alignment = PP_ALIGN.LEFT
		p.level=1
		run = p.add_run()
		run.text = 'Predictive Damage Report for ' + calendar.month_name[int(i.split('_')[1][-2:])]+"-"+i.split('_')[1][:4]
		font = run.font
		font.name = 'Arial'
		font.size = Pt(18.5)
		font.bold = True 
		font.color.rgb = RGBColor(0x00, 0x00, 0x00)
		#Footer for customer report
		left = Inches(6.63)
		top = Inches(6.84)
		width = Inches(3.26)
		height = Inches(0.4)
		s = shapes.add_textbox(left, top, width, height)
		tframe = s.text_frame
		tframe.vertical_anchor = MSO_ANCHOR.MIDDLE
		p = tframe.paragraphs[0]
		p.alignment = PP_ALIGN.RIGHT
		run = p.add_run()
		run.text = 'Customer Report'
		font = run.font
		font.name = 'Arial'
		font.size = Pt(8.2)
		font.color.rgb = RGBColor(127, 127, 127)
		
		# slide = prs.slides.add_slide(prs.slide_masters[1].slide_layouts[5])
		x, y, cx, cy = Inches(0.1), Inches(1.25), Inches(9.8), Inches(5.5)
		fig = plt.figure(figsize=(9.8, 5.5))
		df['ds'] = pd.to_datetime(df['ds'].astype(int)/1000000, unit ='ms')
		df['y'] = df['y'].astype(float).round(2)
		df['yhat_lower'] = (np.where(df['key']=='current', df['y'],df['yhat_lower'])).astype(float).round(2)
		df['yhat_upper'] = (np.where(df['key']=='current', df['y'],df['yhat_upper'])).astype(float).round(2)
		df['month'] = pd.DatetimeIndex(df['ds']).month
		df['month'] = df['month'].apply(lambda x: calendar.month_abbr[x])
		df['err'] = (df['yhat_upper'].astype(float) - df['yhat_lower'].astype(float))/2
		plt.style.use('seaborn')
		plt.plot(range(len(df))[-2:],'yhat_upper',data=df[-2:], c = '#3c9abf',linestyle='-.')
		plt.plot(range(len(df))[-2:],'yhat_lower',data=df[-2:], c = '#3c9abf',linestyle='-.')
		plt.plot(range(len(df)),'y',data=df, c = '#2a6b85',marker='o')
		plt.xticks(range(len(df)), df['month'],fontsize=14)
		plt.fill_between(range(len(df))[-2:], 'yhat_upper', 'yhat_lower',data=df[-2:], color='#3c9abf', alpha='0.5')
		plt.xlabel("Month",weight = 'bold',fontsize=14)
		plt.ylabel("Damages",weight = 'bold',fontsize=14)
		plt.savefig('/var/www/sdmauto/image_stream.png', bbox_inches='tight')
		pic = shapes.add_picture('/var/www/sdmauto/image_stream.png', x, y, cx, cy)

        elif "FAILURE" in str(i):
		##FIRST CHART
		path = "/home/ubuntu/sdmauto/reportdata/" + str(i) +".csv"
		#print path
		df=pd.read_csv(path)
		#name='Consum'
		df = df.loc[df['COMPANYNAME'] == name.encode('utf8')][-12:].reset_index(drop=True)
		#df = df.groupby(['NAME','DATE'], as_index=False).sum()
		#print df.head(100)
		empty_slide = prs.slide_masters[1].slide_layouts[5]
		slide = prs.slides.add_slide(empty_slide)
		shapes = slide.shapes
		left = Inches(0) # 0.93" centers this overall set of shapes
		top = Inches(0.8)
		width = Inches(10)
		height = Inches(0.4)
		# s = shapes3.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
		s = shapes.add_textbox(left, top, width, height)
		fill=s.fill
		fill.solid()
		fill.fore_color.rgb = RGBColor(0xFF, 0xC0, 0x00)
		# s.text ='Step 1'
		tframe = s.text_frame
		# tframe.clear()
		p = tframe.paragraphs[0]
		p.alignment = PP_ALIGN.LEFT
		p.level=1
		run = p.add_run()
		run.text = 'Predictive Failure Report for '+calendar.month_name[int(i.split('_')[1][-2:])]+"-"+i.split('_')[1][:4]
		font = run.font
		font.name = 'Arial'
		font.size = Pt(18.5)
		font.bold = True 
		font.color.rgb = RGBColor(0x00, 0x00, 0x00)
		#Footer for customer report
		left = Inches(6.63)
		top = Inches(6.84)
		width = Inches(3.26)
		height = Inches(0.4)
		s = shapes.add_textbox(left, top, width, height)
		tframe = s.text_frame
		tframe.vertical_anchor = MSO_ANCHOR.MIDDLE
		p = tframe.paragraphs[0]
		p.alignment = PP_ALIGN.RIGHT
		run = p.add_run()
		run.text = 'Customer Report'
		font = run.font
		font.name = 'Arial'
		font.size = Pt(8.2)
		font.color.rgb = RGBColor(127, 127, 127)
		
		# slide = prs.slides.add_slide(prs.slide_masters[1].slide_layouts[5])
		x, y, cx, cy = Inches(0.1), Inches(1.25), Inches(9.8), Inches(5.5)
		fig = plt.figure(figsize=(9.8, 5.5))
		df['ds'] = pd.to_datetime(df['ds'].astype(int)/1000000, unit ='ms')
		df['y'] = df['y'].astype(float).round(2)
		df['yhat_lower'] = (np.where(df['key']=='current', df['y'],df['yhat_lower'])).astype(float).round(2)
		df['yhat_upper'] = (np.where(df['key']=='current', df['y'],df['yhat_upper'])).astype(float).round(2)
		df['month'] = pd.DatetimeIndex(df['ds']).month
		df['month'] = df['month'].apply(lambda x: calendar.month_abbr[x])
		df['err'] = (df['yhat_upper'].astype(float) - df['yhat_lower'].astype(float))/2
		plt.style.use('seaborn')
		plt.plot(range(len(df))[-2:],'yhat_upper',data=df[-2:], c = '#3c9abf',linestyle='-.')
		plt.plot(range(len(df))[-2:],'yhat_lower',data=df[-2:], c = '#3c9abf',linestyle='-.')
		plt.plot(range(len(df)),'y',data=df, c = '#2a6b85',marker='o')
		plt.xticks(range(len(df)), df['month'],fontsize=14)
		plt.fill_between(range(len(df))[-2:], 'yhat_upper', 'yhat_lower',data=df[-2:], color='#3c9abf', alpha='0.5')
		plt.xlabel("Month",weight = 'bold',fontsize=14)
		plt.ylabel("Failure",weight = 'bold',fontsize=14)
		plt.savefig('/var/www/sdmauto/image_stream.png', bbox_inches='tight')
		pic = shapes.add_picture('/var/www/sdmauto/image_stream.png', x, y, cx, cy)

        if "NFF" in str(i):
		##FIRST CHART
		path = "/home/ubuntu/sdmauto/reportdata/" + str(i) +".csv"
		#print path
		df=pd.read_csv(path)
		#name='Consum'
		df = df.loc[df['COMPANYNAME'] == name.encode('utf8')][-12:].reset_index(drop=True)
		#df = df.groupby(['NAME','DATE'], as_index=False).sum()
		#print df.head(100)
		empty_slide = prs.slide_masters[1].slide_layouts[5]
		slide = prs.slides.add_slide(empty_slide)
		shapes = slide.shapes
		left = Inches(0) # 0.93" centers this overall set of shapes
		top = Inches(0.8)
		width = Inches(10)
		height = Inches(0.4)
		# s = shapes3.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
		s = shapes.add_textbox(left, top, width, height)
		fill=s.fill
		fill.solid()
		fill.fore_color.rgb = RGBColor(0xFF, 0xC0, 0x00)
		# s.text ='Step 1'
		tframe = s.text_frame
		# tframe.clear()
		p = tframe.paragraphs[0]
		p.alignment = PP_ALIGN.LEFT
		p.level=1
		run = p.add_run()
		run.text = 'Predictive NFF Report for '+calendar.month_name[int(i.split('_')[1][-2:])]+"-"+i.split('_')[1][:4]
		font = run.font
		font.name = 'Arial'
		font.size = Pt(18.5)
		font.bold = True 
		font.color.rgb = RGBColor(0x00, 0x00, 0x00)
		#Footer for customer report
		left = Inches(6.63)
		top = Inches(6.84)
		width = Inches(3.26)
		height = Inches(0.4)
		s = shapes.add_textbox(left, top, width, height)
		tframe = s.text_frame
		tframe.vertical_anchor = MSO_ANCHOR.MIDDLE
		p = tframe.paragraphs[0]
		p.alignment = PP_ALIGN.RIGHT
		run = p.add_run()
		run.text = 'Customer Report'
		font = run.font
		font.name = 'Arial'
		font.size = Pt(8.2)
		font.color.rgb = RGBColor(127, 127, 127)
		
		# slide = prs.slides.add_slide(prs.slide_masters[1].slide_layouts[5])
		x, y, cx, cy = Inches(0.1), Inches(1.25), Inches(9.8), Inches(5.5)
		fig = plt.figure(figsize=(9.8, 5.5))
		df['ds'] = pd.to_datetime(df['ds'].astype(int)/1000000, unit ='ms')
		df['y'] = df['y'].astype(float).round(2)
		df['yhat_lower'] = (np.where(df['key']=='current', df['y'],df['yhat_lower'])).astype(float).round(2)
		df['yhat_upper'] = (np.where(df['key']=='current', df['y'],df['yhat_upper'])).astype(float).round(2)
		df['month'] = pd.DatetimeIndex(df['ds']).month
		df['month'] = df['month'].apply(lambda x: calendar.month_abbr[x])
		df['err'] = (df['yhat_upper'].astype(float) - df['yhat_lower'].astype(float))/2
		plt.style.use('seaborn')
		plt.plot(range(len(df))[-2:],'yhat_upper',data=df[-2:], c = '#3c9abf',linestyle='-.')
		plt.plot(range(len(df))[-2:],'yhat_lower',data=df[-2:], c = '#3c9abf',linestyle='-.')
		plt.plot(range(len(df)),'y',data=df, c = '#2a6b85',marker='o')
		plt.xticks(range(len(df)), df['month'],fontsize=14)
		plt.fill_between(range(len(df))[-2:], 'yhat_upper', 'yhat_lower',data=df[-2:], color='#3c9abf', alpha='0.5')
		plt.xlabel("Month",weight = 'bold',fontsize=14)
		plt.ylabel("NFF",weight = 'bold',fontsize=14)
		plt.savefig('/var/www/sdmauto/image_stream.png', bbox_inches='tight')
		pic = shapes.add_picture('/var/www/sdmauto/image_stream.png', x, y, cx, cy)


    prs.save('/var/www/sdmauto/test.pptx')
	
    #return render_template_string(completed_template)
  
    try:
        #return "file saved"
        return send_file('/var/www/sdmauto/test.pptx', attachment_filename='/var/www/sdmauto/test.pptx')
    except Exception as e:
        return str(e)
    
if __name__ == '__main__':
    app.run(threaded=True)
#    app.run()

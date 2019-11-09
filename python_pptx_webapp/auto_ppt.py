from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR,PP_ALIGN
from pptx.chart import *
from pptx.enum.chart import XL_CHART_TYPE,XL_LEGEND_POSITION,XL_TICK_MARK,XL_TICK_LABEL_POSITION
from pptx.chart.data import ChartData
import pandas as pd
import numpy as np
from os import listdir
from os.path import isfile, join
import sys
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from StringIO import StringIO

from flask import Flask, render_template_string, request
from wtforms import Form, SelectMultipleField,SelectField
from flask import send_file,send_from_directory


app = Flask(__name__)

##Keep this file as the customer map file. Read it in the UI

usernames=pd.read_csv("/home/ubuntu/sdmauto/Names/DATA_TABLE.csv",error_bad_lines=False,encoding='utf-8')

##Accumulating Reporttype and Date period
mypath = "/home/ubuntu/sdmauto/reportdata"
onlyfiles = pd.DataFrame([f for f in listdir(mypath) if isfile(join(mypath, f))],columns=['Name'])
onlyfiles['Reporttype'] = onlyfiles.Name.str.split('.').str.get(0)


class partnerform(Form):
    #For multiple fields use SelectMultipleFields
    NAME = SelectField(u'Customer Name', choices=[(u, u) for u in sorted(usernames["COMPANY_NAME"].dropna().drop_duplicates().sort_values(),key=unicode.lower)])
    Reporttype = SelectField(u'Report Type ', choices=[(u, u) for u in sorted(onlyfiles["Reporttype"].dropna().drop_duplicates().sort_values())])

template_form = """
{% block content %}
<h1>Select Customer & Report details</h1>
<form method="POST" action="/sdm/report">
    <h2><div>{{ form.NAME.label }}<Label></div></h2>
    <div>{{ form.NAME(multiple=True,size=10) }}</div>
    <h2><div>{{ form.Reporttype.label }}<Label></div></h2>
    <div>{{ form.Reporttype(multiple=True,size=10) }}</div>
    <button type="submit" class="btn">Submit</button>    
</form>
{% endblock %}
"""

completed_template = """
{% block body %}
<h1>Select Customer & Report details</h1>
{% endblock %}
"""

@app.route('/')
def entry():
    form = partnerform(request.form)
    return render_template_string(template_form,form = form)

@app.route('/report',methods=['POST'])
def index():
    #print >>sys.stderr, 'Apache Nonsense!'
    form = partnerform(request.form)
    print "POST request and form is valid"
    name =  form.NAME.data
    print "Partner names %s" % name
    print >>sys.stderr, 'Partner names!'
    ##FIRST SLIDE
    prs = Presentation("/home/ubuntu/data/ppt/ppt.pptx")
    slide = prs.slides[0]
    shapes = slide.shapes    
    shapes[0].text = name.encode('utf8') + ' Monthly Review'
    shapes[1].text = 'Mukesh Choudhary \nJanuary 2018'

    empty_slide = prs.slide_masters[1].slide_layouts[5]
    slide = prs.slides.add_slide(empty_slide)
    shapes = slide.shapes

    ##Grey text box
    left = Inches(0.29) # 0.93" centers this overall set of shapes
    top = Inches(0.58)
    width = Inches(4.55)
    height = Inches(6.92)
    # s = shapes3.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s = shapes.add_textbox(left, top, width, height)
    fill=s.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xEF, 0xF0, 0xF1)

    #Orange textbox with title
    left = Inches(0) # 0.93" centers this overall set of shapes
    top = Inches(2.82)
    width = Inches(10)
    height = Inches(1.82)
    # s = shapes3.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s = shapes.add_textbox(left, top, width, height)
    fill=s.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xC0, 0x00)
    # s.text ='Step 1'
    tframe = s.text_frame
    # tframe.clear()
    tframe.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tframe.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.level=1
    run = p.add_run()
    run.text = 'EXECUTIVE SUMMARY'
    font = run.font
    font.name = 'Arial'
    font.size = Pt(38)
    font.bold = True 
    font.color.rgb = RGBColor(0x00, 0x00, 0x00) 

    #Footer for customer report
    left = Inches(6.63)
    top = Inches(6.84)
    width = Inches(3.26)
    height = Inches(0.4)
    s = shapes.add_textbox(left, top, width, height)
    tframe = s.text_frame
    tframe.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tframe.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = 'Customer Report'
    font = run.font
    font.name = 'Arial'
    font.size = Pt(8.2)
    font.color.rgb = RGBColor(127, 127, 127)
      
    ##FIRST CHART
    path = "/home/ubuntu/sdmauto/reportdata/" + str(form.Reporttype.data) +".csv"
    #print path
    df=pd.read_csv(path)
    #name='Consum'
    df = df.loc[df['COMPANYNAME'] == name.encode('utf8')]
    #df = df.groupby(['NAME','DATE'], as_index=False).sum()
    #print df.head(100)
    empty_slide = prs.slide_masters[1].slide_layouts[5]
    slide = prs.slides.add_slide(empty_slide)
    shapes = slide.shapes
    left = Inches(0) # 0.93" centers this overall set of shapes
    top = Inches(0.8)
    width = Inches(10)
    height = Inches(0.4)
    # s = shapes3.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s = shapes.add_textbox(left, top, width, height)
    fill=s.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xC0, 0x00)
    # s.text ='Step 1'
    tframe = s.text_frame
    # tframe.clear()
    p = tframe.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.level=1
    run = p.add_run()
    run.text = 'Active Devices per Day and Devices on Charge after 10pm TOTAL'
    font = run.font
    font.name = 'Arial'
    font.size = Pt(18.5)
    font.bold = True 
    font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    #Footer for customer report
    left = Inches(6.63)
    top = Inches(6.84)
    width = Inches(3.26)
    height = Inches(0.4)
    s = shapes.add_textbox(left, top, width, height)
    tframe = s.text_frame
    tframe.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tframe.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = 'Customer Report'
    font = run.font
    font.name = 'Arial'
    font.size = Pt(8.2)
    font.color.rgb = RGBColor(127, 127, 127)
 
    # slide = prs.slides.add_slide(prs.slide_masters[1].slide_layouts[5])
    x, y, cx, cy = Inches(0.1), Inches(1.25), Inches(10), Inches(5.5)
    plt.plot('ds','y',data=df)                          

#     image_stream = StringIO()
    plt.savefig('/var/www/sdmauto/image_stream.png')
    pic = shapes.add_picture('/var/www/sdmauto/image_stream.png', x, y, cx, cy)

    prs.save('/var/www/sdmauto/test.pptx')
	
    #return render_template_string(completed_template)
  
    try:
        #return "file saved"
        return send_file('/var/www/sdmauto/test.pptx', attachment_filename='/var/www/sdmauto/test.pptx')
    except Exception as e:
        return str(e)
    
if __name__ == '__main__':
    app.run(threaded=True)
#    app.run()
